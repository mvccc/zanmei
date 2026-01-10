#!/usr/bin/env python3

# vim: set fileencoding=utf-8 :

from datetime import date, timedelta
from pathlib import Path
from pprint import pformat
from typing import Dict, Generator, List, Tuple

import attr
from absl import app, flags, logging as log
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from bible.index import parse_citations
from bible.scripture import BibleVerse, scripture

flags.DEFINE_bool("extract_only", False, "extract text from pptx")
flags.DEFINE_string("pptx", "", "The pptx")
flags.DEFINE_string("master_pptx", "mvccc_master_modern.pptx", "The template pptx")

flags.DEFINE_string("choir", "", "The hymn by choir")
flags.DEFINE_multi_string("hymns", [], "The hymns by congregation")
flags.DEFINE_string("response", "", "The hymn after the teaching")
flags.DEFINE_string("offering", "", "The hymn for the offering")

flags.DEFINE_string("scripture", "", "The bible scriptures")
flags.DEFINE_string("memorize", "", "The bible scripture to memorize")  # verse of the week

flags.DEFINE_string("message", "", "The message")
flags.DEFINE_string("messager", "", "The messager")

flags.DEFINE_bool("communion", None, "Whether to have communion")

FLAGS = flags.FLAGS

PROCESSED = "processed"

# ------------------------------------------------------------------------------


def next_sunday(today: date = None) -> str:
    if today is None:
        today = date.today()
    sunday = today + timedelta(6 - today.weekday())
    return sunday.isoformat()


def extract_slides_text(ppt: Presentation) -> Generator[Tuple[int, List[List[str]]], None, None]:
    for idx, slide in enumerate(ppt.slides):
        shape_text_list: List[List[str]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paragraph_text_list: List[str] = []
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text_list.append("".join(run.text.replace("\xa0", " ").strip() for run in paragraph.runs))
            while paragraph_text_list and not paragraph_text_list[-1]:
                paragraph_text_list.pop()
            if paragraph_text_list:  # Only add non-empty lists
                shape_text_list.append(paragraph_text_list)

        yield idx, shape_text_list


# ------------------------------------------------------------------------------
#
LAYOUT_NAME_PRELUDE = ("prelude",)
LAYOUT_NAME_MESSAGE = ("message",)
LAYOUT_NAME_HYMN = ("hymn",)
LAYOUT_NAME_SCRIPTURE = ("verse",)
LAYOUT_NAME_MEMORIZE = ("memorize",)
LAYOUT_NAME_TEACHING = ("teaching",)
LAYOUT_NAME_SECTION = ("section",)
LAYOUT_NAME_BLANK = ("blank",)


def _get_slide_layout(ppt: Presentation, layout_names: Tuple[str, ...]):
    normalized = {n.strip().lower() for n in layout_names}
    for layout in ppt.slide_layouts:
        if layout.name and layout.name.strip().lower() in normalized:
            return layout
    available = [layout.name for layout in ppt.slide_layouts]
    raise ValueError(f"Could not find slide layout in {layout_names}. Available layouts: {available}")


def _get_placeholder_by_type(container, placeholder_types: Tuple[PP_PLACEHOLDER, ...]):
    for placeholder in container.placeholders:
        pf = getattr(placeholder, "placeholder_format", None)
        ptype = getattr(pf, "type", None)
        if ptype is None:
            continue
        if any(ptype == t for t in placeholder_types):
            return placeholder
    available = []
    for placeholder in container.placeholders:
        pf = getattr(placeholder, "placeholder_format", None)
        ptype = getattr(pf, "type", None)
        if ptype is None:
            continue
        available.append(ptype)
    raise ValueError(f"Could not find placeholder types {placeholder_types}. Available placeholder types: {available}")


@attr.s
class Prelude:
    message: str = attr.ib()
    picture: str = attr.ib()

    def add_to(self, ppt: Presentation, padding="\u3000\u3000") -> Presentation:
        slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_PRELUDE))
        # Try TITLE first, fall back to BODY if not available
        try:
            title = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
        except ValueError:
            title = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.BODY,))
        picture = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT))

        title.text = padding + self.message
        picture.insert_picture(self.picture)

        return ppt


@attr.s
class Message:
    message: str = attr.ib()

    def add_to(self, ppt: Presentation, padding="\u3000\u3000") -> Presentation:
        slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_MESSAGE))
        body = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.BODY,))
        body.text = padding + self.message

        return ppt


@attr.s
class Hymn:
    filename: str = attr.ib()  # can be index number, hymn's title
    lyrics: List[Tuple[str, List[str]]] = attr.ib()  # List[title, paragraph]

    def add_to(self, ppt: Presentation, padding: str = " ") -> Presentation:
        for _, slide_content in self.lyrics:
            # Handle both cases: [title_list, paragraph_list] or [combined_list]
            if len(slide_content) == 2:
                title, paragraph = slide_content
            elif len(slide_content) == 1:
                # If only one list, treat it as paragraph with empty title
                title = [""]
                paragraph = slide_content[0]
            else:
                raise ValueError(f"Unexpected slide_content structure: {slide_content}")

            slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_HYMN))
            title_holder = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
            paragraph_holder = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.BODY,))
            title_holder.text = title[0]
            # XXX: workaround alignment problem
            if paragraph:
                paragraph[0] = padding + paragraph[0]
            paragraph_holder.text = "\n".join(paragraph)

        return ppt


def search_hymn_ppt(keyword: str, basepath: Path = None) -> List[Hymn]:
    if basepath is None:
        basepath = Path(PROCESSED)

    keyword = keyword.replace(".pptx", "")
    ptn = f"**/*{keyword}*.pptx"
    glob = basepath.glob(ptn)
    found = list(glob)

    if not found:
        # interchangeability characters
        interchangebles = [("你", "祢", "袮"), ("寶", "寳"), ("他", "祂"), ("于", "於"), ("牆", "墻")]
        for t in interchangebles:
            for w in t:
                if w not in ptn:
                    continue
                for w1 in t:
                    if w == w1:
                        continue
                    glob = basepath.glob(ptn.replace(w, w1))
                    found = list(glob)
                    if found:
                        break

    if not found:
        # Create a placeholder hymn with title and empty slide
        log.warning(f"can not find anything match {ptn}. Creating placeholder slides.")
        # Create a simple hymn structure: title slide + empty slide
        placeholder_hymn = Hymn(filename=keyword, lyrics=[(0, [[keyword], ["(歌詞待補充)"]]), (1, [[""], [""]])])
        return [placeholder_hymn]

    if len(found) > 1:
        log.warning(f"found more than 1 files for {ptn}. {[p.as_posix() for p in found]}")

    found = [path for path in found if path.stem == keyword] + [path for path in found if path.stem != keyword]

    result: List[Hymn] = []
    for path in found:
        ppt = Presentation(path.as_posix())
        lyrics = list(extract_slides_text(ppt))
        hymn = Hymn(path.name, lyrics)
        log.info(f"keyword={keyword}, lyrics=\n{pformat(hymn.lyrics)}")
        result.append(hymn)

    return result


@attr.s
class Section:
    title: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_SECTION))
        # Try TITLE first, fall back to BODY if not available
        try:
            title_ph = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
        except ValueError:
            title_ph = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.BODY,))
        title_ph.text = self.title

        return ppt


@attr.s
class Scripture:
    citations: str = attr.ib()
    cite_verses: Dict[str, List[BibleVerse]] = attr.ib()

    def add_to(self, ppt: Presentation, padding="  ") -> Presentation:
        for cite, verses in self.cite_verses.items():
            for idx, bv in enumerate(verses):
                if idx % 2 == 0:
                    slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_SCRIPTURE))
                    title_ph = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
                    message = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.BODY,))
                    title_ph.text = cite
                    message.text = ""
                message.text += (padding if idx % 2 == 0 else "\n") + f"{bv.verse}\u3000{bv.text}"

        return ppt


def to_scripture(citations: str) -> Scripture:
    bible = scripture()
    cite_verses = bible.search(parse_citations(citations).items())
    for cite, verses in cite_verses.items():
        log.info(f"citation={cite}, verses=\n{pformat(verses)}")

    return Scripture(citations, cite_verses)


@attr.s
class Memorize:
    citation: str = attr.ib()
    verses: List[BibleVerse] = attr.ib()

    def add_to(self, ppt: Presentation, padding="\u3000\u3000") -> Presentation:
        slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_MEMORIZE))
        title_ph = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
        message = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.BODY,))
        title_ph.text = "本週金句"
        message.text = padding + "".join(bv.text for bv in self.verses) + f"\n\n{self.citation:>35}"

        return ppt


@attr.s
class Teaching:
    title: str = attr.ib()
    message: str = attr.ib()
    messenger: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_TEACHING))
        title_ph = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
        # Try to get BODY, if not available, put all content in title
        try:
            body = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.BODY,))
            title_ph.text = self.title
            body.text = f"{self.message}\n\n{self.messenger}"
        except ValueError:
            # No BODY placeholder, combine all in title
            title_ph.text = f"{self.title}\n{self.message}\n{self.messenger}"

        return ppt


@attr.s
class Blank:
    def add_to(self, ppt: Presentation) -> Presentation:
        _ = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_BLANK))

        return ppt


def mvccc_slides(
    hymns: List[str],
    scripture: str,
    memorize: str,
    message: str,
    messager: str,
    choir: str,
    response: str,
    offering: str,
    communion: bool,
) -> List:
    slides = [
        Message("請儘量往前或往中間坐,並將手機關閉或關至靜音,預備心敬拜！"),
        Message(
            """惟耶和華在他的聖殿中；全地的人，都當在他面前肅敬靜默。

                    哈巴谷書 2:20"""
        ),
    ]
    hymn = search_hymn_ppt("聖哉聖哉聖哉")
    slides.append(hymn[0])

    slides.append(Section("宣  召"))

    slides.append(Section("頌  讚"))
    for kw in hymns:
        r = search_hymn_ppt(kw)
        slides.append(r[0])

    slides.append(Section("祈  禱"))

    slides.append(Section("讀  經"))

    slides.append(to_scripture(scripture))
    for cite, verses in to_scripture(memorize).cite_verses.items():
        slides.append(Memorize(cite, verses))
        break
    slides.append(Blank())

    slides.append(Section("獻  詩"))
    if choir:
        hymn = search_hymn_ppt(choir)[0]
        slides.append(hymn)

    slides.append(Teaching("信息", f"「{message}」", f"{messager}"))

    slides.append(Section("回  應"))
    if response:
        hymn = search_hymn_ppt(response)[0]
        slides.append(hymn)

    if offering:
        hymn = search_hymn_ppt(offering)[0]
        slides.append(hymn)

    slides.append(Section("奉 獻 禱 告"))

    if communion:
        slides.append(Section("聖  餐"))

    slides.append(Section("歡 迎 您"))
    slides.append(Section("家 事 分 享"))

    hymn = search_hymn_ppt("三一頌")[0]
    slides.append(hymn)

    slides.append(Section("祝  福"))
    slides.append(Section("默  禱"))
    slides.append(Blank())

    return slides


def to_pptx(slides: List, master_slide: Presentation) -> Presentation:
    # Create a new presentation using the master template's layouts
    # but without any of the master's existing slides
    ppt = Presentation(FLAGS.master_pptx)

    # Remove all existing slides from the template
    while len(ppt.slides) > 0:
        r_id = ppt.slides._sldIdLst[0].rId
        ppt.part.drop_rel(r_id)
        del ppt.slides._sldIdLst[0]

    for slide in slides:
        slide.add_to(ppt)

    return ppt


# ------------------------------------------------------------------------------


def main(argv):
    del argv

    if FLAGS.extract_only:
        ppt = Presentation(FLAGS.pptx)

        for _idx, _text in extract_slides_text(ppt):
            pass
        return

    slides = mvccc_slides(
        hymns=FLAGS.hymns,
        scripture=FLAGS.scripture,
        memorize=FLAGS.memorize,
        message=FLAGS.message,
        messager=FLAGS.messager,
        choir=FLAGS.choir,
        response=FLAGS.response,
        offering=FLAGS.offering,
        communion=FLAGS.communion,
    )
    master = Presentation(FLAGS.master_pptx)
    ppt = to_pptx(slides, master)
    ppt.save(FLAGS.pptx)


if __name__ == "__main__":
    app.run(main)
