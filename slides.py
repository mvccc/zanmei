#!/usr/bin/env python3

# vim: set fileencoding=utf-8 :

from datetime import date, timedelta
from pathlib import Path
from pprint import pformat
from typing import Generator, List, Tuple

import attr
from absl import app, flags, logging as log

import scripture
from pptx import Presentation

flags.DEFINE_bool("extract_only", False, "extract text from pptx")
flags.DEFINE_string("pptx", "", "The pptx")
flags.DEFINE_string("master_pptx", "mvccc_master.pptx", "The template pptx")

flags.DEFINE_string("choir", "", "The hymn by choir")
flags.DEFINE_multi_string("hymns", [], "The hymns by congregation")
flags.DEFINE_string("response", "", "The hymn after the teaching")
flags.DEFINE_string("offering", "", "The hymn for the offering")

flags.DEFINE_string("scripture", "", "The bible scriptures")
flags.DEFINE_string("memorize", "", "The bible scripture to memorize")  # verse of the week

flags.DEFINE_string("message", "", "The message")
flags.DEFINE_string("messager", "", "The messager")

flags.DEFINE_bool("communion", False, "Whether to have communion")

FLAGS = flags.FLAGS

PROCESSED = Path("processed")

# ------------------------------------------------------------------------------


def next_sunday() -> str:
    today = date.today()
    sunday = today + timedelta(6 - today.weekday())
    return sunday.isoformat()


def extract_slides_text(ppt: Presentation) -> Generator[Tuple[int, List[List[str]]], None, None]:
    for idx, slide in enumerate(ppt.slides):
        shape_text_list: List[List[str]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paragraph_text_list: List[str] = []
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text_list.append("".join(run.text.replace("\xa0", " ") for run in paragraph.runs))
            while not paragraph_text_list[-1]:
                paragraph_text_list.pop()
            shape_text_list.append(paragraph_text_list)

        yield idx, shape_text_list


# ------------------------------------------------------------------------------

LAYOUT_PRELUDE = 0
LAYOUT_MESSAGE = 1
LAYOUT_HYMN = 2
LAYOUT_VERSE = 3
LAYOUT_MEMORIZE = 4
LAYOUT_TEACHING = 5
LAYOUT_SECTION = 6
LAYOUT_BLANK = 7


@attr.s
class Prelude:
    message: str = attr.ib()
    picture: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_PRELUDE])
        message, picture = slide.placeholders
        message.text = self.message
        picture.insert_picture(self.picture)

        return ppt


@attr.s
class Message:
    message: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_MESSAGE])
        message, = slide.placeholders
        message.text = self.message

        return ppt


@attr.s
class Hymn:
    index: str = attr.ib()  # can be index number, hymn's title
    lyrics: List[Tuple[str, List[str]]] = attr.ib(init=False)  # List[title, paragraph]

    def __attrs_post_init__(self):
        ptn = f"**/*{self.index}*.pptx"
        glob = PROCESSED.glob(ptn)
        found = list(glob)
        assert found, f"can not find anything match {ptn}."
        if len(found) > 1:
            log.warn(f"found more than 1 files for {ptn}. {[p.as_posix() for p in found]}")
        for path in found:
            if path.stem == self.index:
                filepath = path.as_posix()
                break
        else:
            filepath = found[0].as_posix()
        ppt = Presentation(filepath)
        self.lyrics = list(extract_slides_text(ppt))
        log.info(f"index={self.index}, lyrics=\n{pformat(self.lyrics)}")

    def add_to(self, ppt: Presentation) -> Presentation:
        for _, (title, paragraph) in self.lyrics:
            slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_HYMN])
            title_holder, paragraph_holder = slide.placeholders
            title_holder.text = title[0]
            # XXX: workaround alignment problem
            paragraph[0] = " " + paragraph[0]
            paragraph_holder.text = "\n".join(paragraph)

        return ppt


@attr.s
class Section:
    title: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_SECTION])
        title, = slide.placeholders
        title.text = self.title

        return ppt


@attr.s
class Scripture:
    locations: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        bible = scripture.scripture()
        for loc, verses in scripture.search(bible, self.locations).items():
            for t in verses.itertuples():
                slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_MEMORIZE])
                title, message = slide.placeholders
                title.text = loc
                message.text = f"{t.Index[-1]} {t.scripture}"  # verse: scripture

        return ppt


@attr.s
class Memorize:
    location: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_MEMORIZE])
        title, message = slide.placeholders

        title.text = "本週金句"

        bible = scripture.scripture()
        loc, verses = list(scripture.search(bible, self.location).items())[0]
        text = "".join(verses["scripture"].to_list())
        message.text = text + f"\n{loc:>20}"

        return ppt


@attr.s
class Teaching:
    # TODO
    title: str = attr.ib()
    message: str = attr.ib()
    messenger: str = attr.ib()

    def add_to(self, ppt: Presentation) -> Presentation:
        slide = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_TEACHING])
        message, = slide.placeholders
        message.text = "\n\n".join([self.title, self.message, self.messenger])

        return ppt


@attr.s
class Blank:
    def add_to(self, ppt: Presentation) -> Presentation:
        _ = ppt.slides.add_slide(ppt.slide_layouts[LAYOUT_BLANK])

        return ppt


def mvccc_slides() -> List:
    slides = [
        Prelude("  請儘量往前或往中間坐,並將手機關閉或關至靜音,預備心敬拜！", "silence_phone1.png"),
        Message(
            """  惟耶和華在他的聖殿中；
全地的人，都當在他面前肅敬靜默。

            哈巴谷書 2:20"""
        ),
    ]
    slides.append(Hymn("聖哉聖哉聖哉"))

    slides.append(Section("頌  讚"))
    if FLAGS.choir:
        slides.append(Hymn(FLAGS.choir))

    slides.append(Section("宣  召"))
    slides.append(Section("祈  禱"))

    slides.append(Section("讀  經"))
    slides.append(Scripture(FLAGS.scripture))
    slides.append(Memorize(FLAGS.memorize))

    slides.append(Section("獻  詩"))
    if FLAGS.hymns:
        slides.extend(map(Hymn, FLAGS.hymns))

    slides.append(Teaching("信息", f"「{FLAGS.message}」", f"{FLAGS.messager}"))

    slides.append(Section("回  應"))
    if FLAGS.response:
        slides.append(Hymn(FLAGS.response))

    if FLAGS.offering:
        slides.append(Hymn(FLAGS.offering))

    slides.append(Section("奉 獻 禱 告"))
    slides.append(Section("歡 迎 您"))
    slides.append(Section("家 事 分 享"))

    slides.append(Hymn("三一頌"))

    if FLAGS.communion:
        slides.append(Section("聖  餐"))

    slides.append(Section("祝  福"))
    slides.append(Section("默  禱"))
    slides.append(Blank())

    return slides


def to_pptx(slides: List, master_slide: Presentation) -> Presentation:
    ppt = master_slide

    for slide in slides:
        slide.add_to(ppt)

    return ppt


# ------------------------------------------------------------------------------


def main(argv):
    del argv

    if FLAGS.extract_only:
        ppt = Presentation(FLAGS.pptx)

        for idx, text in extract_slides_text(ppt):
            print(f"{idx+1:02d} {text}")
        return

    slides = mvccc_slides()
    master = Presentation(FLAGS.master_pptx)
    ppt = to_pptx(slides, master)
    ppt.save(FLAGS.pptx)


if __name__ == "__main__":
    app.run(main)
