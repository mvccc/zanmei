#!/usr/bin/env python3
"""Convert hymn lyrics between PPTX and Markdown formats.

Markdown format:
    # Hymn Title

    ## (1)
    First line of verse 1
    Second line of verse 1

    ## (2)
    First line of verse 2
    Second line of verse 2
"""

import re
from pathlib import Path

from pptx import Presentation


def pptx_to_md(pptx_path: str | Path) -> str:
    from mvccc.slides import extract_slides_text

    ppt = Presentation(str(pptx_path))
    slides = list(extract_slides_text(ppt))

    if not slides:
        return ""

    lines = []
    hymn_title = None

    for _idx, slide_content in slides:
        if len(slide_content) == 1:
            title_list = slide_content[0]
            title = title_list[0] if title_list else ""
            if hymn_title is None:
                hymn_title = _clean_title(title)
                lines.append(f"# {hymn_title}")
                lines.append("")
            else:
                lines.append("##")
                for lyric_line in title_list:
                    if lyric_line.strip():
                        lines.append(_format_lyric_line(lyric_line))
                lines.append("")
            continue

        if len(slide_content) >= 2:
            title_list, lyrics_list = slide_content[0], slide_content[1]
            title = title_list[0] if title_list else ""

            if hymn_title is None:
                hymn_title = _clean_title(title)
                lines.append(f"# {hymn_title}")
                lines.append("")

            verse_marker = _extract_verse_marker(title)
            if verse_marker:
                lines.append(f"## {verse_marker}")
            else:
                lines.append("##")

            for lyric_line in lyrics_list:
                lines.append(_format_lyric_line(lyric_line))
            lines.append("")

    return "\n".join(lines).rstrip() + "\n"


def _clean_title(title: str) -> str:
    title = re.sub(r"^#?\d+[_\s]*", "", title)
    title = re.sub(r"\s*\(\d+\)\s*$", "", title)
    return title.strip()


def _extract_verse_marker(title: str) -> str:
    match = re.search(r"\((\d+)\)\s*$", title)
    if match:
        return f"({match.group(1)})"
    return ""


def _format_lyric_line(line: str) -> str:
    if not line.strip():
        return ""
    if line.endswith("  "):
        return line
    return f"{line.rstrip()}  "


def _strip_md_linebreak(line: str) -> str:
    if not line.strip():
        return ""
    if line.endswith("  "):
        return line[:-2]
    return line.rstrip()


def md_to_lyrics(md_content: str) -> tuple[str, list[tuple[str, list[str]]]]:
    """Returns: (hymn_title, [(verse_marker, [lyric_lines]), ...])"""
    lines = md_content.strip().split("\n")

    hymn_title = ""
    slides = []
    current_verse = ""
    current_lyrics = []

    for line in lines:
        if line.startswith("# "):
            hymn_title = line[2:].strip()
        elif line.startswith("## "):
            if current_lyrics:
                slides.append((current_verse, current_lyrics))
            current_verse = line[3:].strip()
            current_lyrics = []
        elif line.startswith("##"):
            if current_lyrics:
                slides.append((current_verse, current_lyrics))
            current_verse = ""
            current_lyrics = []
        elif line.strip():
            current_lyrics.append(_strip_md_linebreak(line))

    if current_lyrics:
        slides.append((current_verse, current_lyrics))

    return hymn_title, slides


def split_long_lines(lyrics: list[str], min_length: int = 30) -> list[str]:
    result = []
    sentence_delims = "；!?！？。?"

    for line in lyrics:
        if len(line) < min_length:
            result.append(line)
            continue

        if "；" in line:
            parts = [p.strip() for p in line.split("；") if p.strip()]
            if len(parts) >= 2:
                for i, part in enumerate(parts):
                    if i < len(parts) - 1:
                        result.append(part + "；")
                    else:
                        result.append(part)
                continue

        segments = []
        buffer = ""
        for char in line:
            buffer += char
            if char in sentence_delims:
                segments.append(buffer.strip())
                buffer = ""
        if buffer.strip():
            segments.append(buffer.strip())

        if len(segments) <= 1:
            result.append(line)
            continue

        for segment in segments:
            if len(segment) >= min_length and "，" in segment:
                parts = [p.strip() for p in segment.split("，") if p.strip()]
                if len(parts) >= 2:
                    for i, part in enumerate(parts):
                        if i < len(parts) - 1:
                            result.append(part + "，")
                        else:
                            result.append(part)
                    continue
            result.append(segment)

    return result


def reformat_lyrics(
    slides: list[tuple[str, list[str]]], max_lines: int = 5, target_lines: int = 4
) -> list[tuple[str, list[str]]]:
    """Split slides >max_lines into target_lines chunks. Allows 5 lines for trailing "阿們"."""
    result = []

    for verse_marker, lyrics in slides:
        lyrics = split_long_lines(lyrics)

        if len(lyrics) <= max_lines:
            result.append((verse_marker, lyrics))
            continue

        chunks = []
        for i in range(0, len(lyrics), target_lines):
            chunk = lyrics[i : i + target_lines]
            remaining = len(lyrics) - (i + target_lines)
            if 0 < remaining <= (max_lines - target_lines):
                chunks.append(lyrics[i:])
                break
            chunks.append(chunk)

        for chunk in chunks:
            result.append((verse_marker, chunk))

    return result


def lyrics_to_md(hymn_title: str, slides: list[tuple[str, list[str]]]) -> str:
    lines = [f"# {hymn_title}", ""]

    for verse_marker, lyrics in slides:
        if verse_marker:
            lines.append(f"## {verse_marker}")
        else:
            lines.append("##")
        lines.extend(_format_lyric_line(lyric) for lyric in lyrics)
        lines.append("")

    return "\n".join(lines).rstrip() + "\n"


def reformat_md(md_content: str, max_lines: int = 5, target_lines: int = 4) -> str:
    hymn_title, slides = md_to_lyrics(md_content)
    reformatted_slides = reformat_lyrics(slides, max_lines, target_lines)
    return lyrics_to_md(hymn_title, reformatted_slides)


def md_to_pptx(md_path: str | Path, template_pptx: str = "mvccc_master_modern_dark.pptx") -> Presentation:
    from pptx.enum.shapes import PP_PLACEHOLDER

    from mvccc.slides import (
        LAYOUT_NAME_HYMN_LYRICS,
        LAYOUT_NAME_HYMN_TITLE,
        _get_placeholder_by_type,
        _get_slide_layout,
    )

    md_path = Path(md_path)
    md_content = md_path.read_text(encoding="utf-8")
    hymn_title, slides_data = md_to_lyrics(md_content)

    ppt = Presentation(template_pptx)

    while len(ppt.slides) > 0:
        r_id = ppt.slides._sldIdLst[0].rId
        ppt.part.drop_rel(r_id)
        del ppt.slides._sldIdLst[0]

    title_slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_HYMN_TITLE))
    title_holder = _get_placeholder_by_type(title_slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
    title_holder.text = hymn_title
    try:
        body = _get_placeholder_by_type(title_slide, (PP_PLACEHOLDER.BODY,))
        body.text = ""
    except ValueError:
        pass

    for _verse_marker, lyrics in slides_data:
        slide = ppt.slides.add_slide(_get_slide_layout(ppt, LAYOUT_NAME_HYMN_LYRICS))
        body = _get_placeholder_by_type(slide, (PP_PLACEHOLDER.BODY,))
        body.text = " " + "\n".join(lyrics)

    return ppt


def convert_pptx_to_md(pptx_path: str | Path) -> Path:
    pptx_path = Path(pptx_path)
    md_content = pptx_to_md(pptx_path)
    md_path = pptx_path.with_suffix(".md")
    md_path.write_text(md_content, encoding="utf-8")
    return md_path


def convert_md_to_pptx(md_path: str | Path, template_pptx: str = "mvccc_master_modern_dark.pptx") -> Path:
    md_path = Path(md_path)
    ppt = md_to_pptx(md_path, template_pptx)
    pptx_path = md_path.with_suffix(".pptx")
    ppt.save(str(pptx_path))
    return pptx_path


def convert_all_pptx(directory: str | Path = "processed/mvccc") -> list[Path]:
    directory = Path(directory)
    converted = []

    for pptx_path in sorted(directory.glob("*.pptx")):
        try:
            md_path = convert_pptx_to_md(pptx_path)
            converted.append(md_path)
        except Exception as e:
            print(f"Error converting {pptx_path}: {e}")  # noqa: T201

    return converted


def reformat_md_file(md_path: str | Path) -> Path:
    md_path = Path(md_path)
    content = md_path.read_text(encoding="utf-8")
    reformatted = reformat_md(content)
    md_path.write_text(reformatted, encoding="utf-8")
    return md_path


def reformat_all_md(directory: str | Path = "processed/mvccc") -> list[Path]:
    directory = Path(directory)
    reformatted = []

    for md_path in sorted(directory.glob("*.md")):
        try:
            content = md_path.read_text(encoding="utf-8")
            new_content = reformat_md(content)
            if content != new_content:
                md_path.write_text(new_content, encoding="utf-8")
                reformatted.append(md_path)
        except Exception as e:
            print(f"Error reformatting {md_path}: {e}")  # noqa: T201

    return reformatted


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage:")  # noqa: T201
        print("  python -m hymns.hymn_md <file.pptx>       Convert PPTX to MD")  # noqa: T201
        print("  python -m hymns.hymn_md <file.md>         Convert MD to PPTX")  # noqa: T201
        print("  python -m hymns.hymn_md --all             Convert all PPTX in processed/mvccc to MD")  # noqa: T201
        print("  python -m hymns.hymn_md --reformat        Reformat all MD files (split >5 lines)")  # noqa: T201
        print("  python -m hymns.hymn_md --reformat <file> Reformat single MD file")  # noqa: T201
        sys.exit(1)

    if sys.argv[1] == "--all":
        converted = convert_all_pptx()
        print(f"Converted {len(converted)} files")  # noqa: T201
    elif sys.argv[1] == "--reformat":
        if len(sys.argv) > 2:
            md_path = reformat_md_file(sys.argv[2])
            print(f"Reformatted: {md_path}")  # noqa: T201
        else:
            reformatted = reformat_all_md()
            print(f"Reformatted {len(reformatted)} files")  # noqa: T201
    elif sys.argv[1].endswith(".pptx"):
        md_path = convert_pptx_to_md(sys.argv[1])
        print(f"Converted: {md_path}")  # noqa: T201
    elif sys.argv[1].endswith(".md"):
        pptx_path = convert_md_to_pptx(sys.argv[1])
        print(f"Converted: {pptx_path}")  # noqa: T201
    else:
        print(f"Unknown file type: {sys.argv[1]}")  # noqa: T201
        sys.exit(1)
