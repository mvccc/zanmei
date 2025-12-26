#!/usr/bin/env python3
"""
Create sample slides with modern design that can be used as reference
for updating the master template.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Modern color palette
COLORS = {
    'background': RGBColor(245, 247, 250),     # Light grayish-blue #F5F7FA
    'primary': RGBColor(41, 98, 255),          # Modern blue #2962FF
    'secondary': RGBColor(100, 116, 139),      # Slate gray #64748B
    'text': RGBColor(30, 41, 59),              # Dark slate #1E293B
    'accent': RGBColor(139, 92, 246),          # Purple #8B5CF6
    'white': RGBColor(255, 255, 255),          # White
}

def create_sample_slides():
    """Create sample slides with modern styling"""
    ppt = Presentation()
    ppt.slide_width = Inches(13.33)
    ppt.slide_height = Inches(7.5)
    
    # Sample 1: Title/Section slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['background']
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.33), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "頌  讚"
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "PingFang TC"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Sample 2: Content slide (Message/Hymn)
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['background']
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(11.33), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "讚美真神萬福之根"
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "PingFang TC"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Body text
    body_box = slide.shapes.add_textbox(
        Inches(2), Inches(2.5), Inches(9.33), Inches(4)
    )
    body_frame = body_box.text_frame
    body_frame.text = """讚美真神萬福之根
世上萬民讚美主恩
天使天軍讚美主名
讚美聖父聖子聖靈
阿們"""
    
    body_frame.word_wrap = True
    for paragraph in body_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = "PingFang TC"
        paragraph.font.size = Pt(32)
        paragraph.font.color.rgb = COLORS['text']
        paragraph.space_before = Pt(12)
    
    # Sample 3: Scripture slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['background']
    
    # Citation
    citation_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(11.33), Inches(0.8)
    )
    citation_frame = citation_box.text_frame
    citation_frame.text = "羅馬書 3:23-24"
    
    p = citation_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "PingFang TC"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Verse text
    verse_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2), Inches(10.33), Inches(4)
    )
    verse_frame = verse_box.text_frame
    verse_frame.text = "因為世人都犯了罪，虧缺了　神的榮耀；如今卻蒙　神的恩典，因基督耶穌的救贖，就白白地稱義。"
    
    verse_frame.word_wrap = True
    for paragraph in verse_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = "PingFang TC"
        paragraph.font.size = Pt(36)
        paragraph.font.color.rgb = COLORS['text']
        paragraph.line_spacing = 1.5
    
    ppt.save('modern_design_samples.pptx')
    print("✓ Created modern_design_samples.pptx with 3 sample slides")
    print("\nModern Design Specifications:")
    print("=" * 60)
    print(f"Background Color: RGB{COLORS['background']}")
    print(f"Primary Color (Titles): RGB{COLORS['primary']}")
    print(f"Text Color: RGB{COLORS['text']}")
    print(f"Title Font: PingFang TC (macOS) or Microsoft YaHei (Windows)")
    print(f"Body Font: PingFang TC (macOS) or Microsoft YaHei (Windows)")
    print(f"Title Size: 48-72pt")
    print(f"Body Size: 32-36pt")
    print("\nTo create modern_master.pptx:")
    print("1. Open mvccc_master.pptx in PowerPoint or Keynote")
    print("2. Go to View > Edit Master Slides")
    print("3. Apply the colors and fonts from modern_design_samples.pptx")
    print("4. Save as modern_master.pptx")

if __name__ == "__main__":
    create_sample_slides()
